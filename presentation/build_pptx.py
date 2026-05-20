"""
Nex Order — Pitch deck generator.

Run:
    pip install python-pptx
    python presentation/build_pptx.py

Produces: presentation/NexOrder-Pitch.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


OUTPUT_PATH = Path(__file__).parent / "NexOrder-Pitch.pptx"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

INK = RGBColor(0x1F, 0x29, 0x37)        # deep slate body / headers
MUTED = RGBColor(0x6B, 0x72, 0x80)       # secondary copy
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)      # cyan accent
HAIRLINE = RGBColor(0xE5, 0xE7, 0xEB)    # dividers / outlines
SOFT_BG = RGBColor(0xF8, 0xFA, 0xFC)     # block backgrounds
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_DISPLAY = "Calibri"
FONT_BODY = "Calibri"


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def set_solid_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_no_fill(shape) -> None:
    shape.fill.background()


def set_line(shape, rgb: RGBColor | None, width_pt: float = 0.75) -> None:
    if rgb is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    return tb, tf


def style_run(run, *, size_pt: float, bold: bool = False, color: RGBColor = INK,
              font: str = FONT_BODY) -> None:
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_paragraph(tf, text: str, *, size_pt: float, bold: bool = False,
                  color: RGBColor = INK, font: str = FONT_BODY,
                  align=PP_ALIGN.LEFT, space_after: float = 4) -> None:
    p = tf.add_paragraph() if tf.paragraphs[0].text or tf.paragraphs[0].runs else tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    style_run(run, size_pt=size_pt, bold=bold, color=color, font=font)


def add_accent_bar(slide, *, left, top, width=Inches(0.6), height=Inches(0.08)) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_solid_fill(bar, ACCENT)
    set_line(bar, None)


def add_speaker_notes(slide, text: str) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text
    for p in notes_tf.paragraphs:
        for r in p.runs:
            style_run(r, size_pt=11, color=INK)


def add_footer(slide, text: str) -> None:
    tb, tf = add_textbox(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
    add_paragraph(tf, text, size_pt=9, color=MUTED)
    # hairline above
    line = slide.shapes.add_connector(1, Inches(0.6), Inches(7.0),
                                      Inches(12.7), Inches(7.0))
    line.line.color.rgb = HAIRLINE
    line.line.width = Pt(0.5)


# ---------------------------------------------------------------------------
# Slide 1 — Title / Positioning
# ---------------------------------------------------------------------------

def build_slide_1(prs: Presentation) -> None:
    slide = add_blank_slide(prs)

    # Top-left wordmark band
    add_accent_bar(slide, left=Inches(0.6), top=Inches(0.7))

    tb, tf = add_textbox(slide, Inches(0.6), Inches(0.85), Inches(6), Inches(0.4))
    add_paragraph(tf, "NEX ORDER", size_pt=12, bold=True, color=ACCENT,
                  font=FONT_DISPLAY)

    # Headline
    tb, tf = add_textbox(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(2.5))
    add_paragraph(tf, "One pipeline for every B2B order.",
                  size_pt=48, bold=True, color=INK, font=FONT_DISPLAY,
                  space_after=14)
    add_paragraph(tf,
                  "Multi-channel order intake, tier-aware pricing, "
                  "route-driven field sales, and a live audit trail. "
                  "Built for distributors and wholesalers.",
                  size_pt=20, color=MUTED, space_after=0)

    # Bottom proof strip
    tb, tf = add_textbox(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(1.0))
    add_paragraph(tf,
                  "Proven on a live distributor deployment "
                  "running on Supabase + Vercel.",
                  size_pt=14, color=INK, space_after=4)
    add_paragraph(tf, "Live demo: nexorder.vercel.app",
                  size_pt=12, color=ACCENT, bold=True)

    add_footer(slide, "Nex Order   •   Pitch deck   •   Page 1 of 5")

    add_speaker_notes(slide,
        "Open with positioning. Nex Order is a generic B2B order-management "
        "platform that unifies every channel a distributor sells through — "
        "self-serve customers, telesales, field reps, and future EDI partners. "
        "What makes it different in 2026 isn't the UI, it's the architecture: "
        "one server-side gate makes pricing, promotions, stock, and audit "
        "behave identically no matter where the order came from. We'll spend "
        "10 minutes on slides, then move into a live walkthrough on the URL "
        "shown at the bottom of this slide.")


# ---------------------------------------------------------------------------
# Slide 2 — The Problem
# ---------------------------------------------------------------------------

def build_slide_2(prs: Presentation) -> None:
    slide = add_blank_slide(prs)

    # Eyebrow
    add_accent_bar(slide, left=Inches(0.6), top=Inches(0.6))
    tb, tf = add_textbox(slide, Inches(0.6), Inches(0.78), Inches(6), Inches(0.35))
    add_paragraph(tf, "THE PROBLEM", size_pt=11, bold=True, color=ACCENT,
                  font=FONT_DISPLAY)

    # Headline
    tb, tf = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.4))
    add_paragraph(tf,
                  "Orders come from everywhere.",
                  size_pt=34, bold=True, color=INK, font=FONT_DISPLAY,
                  space_after=4)
    add_paragraph(tf,
                  "Most ops teams can't see them in one place.",
                  size_pt=24, color=MUTED, space_after=0)

    # Three pain blocks
    block_top = Inches(3.4)
    block_h = Inches(2.6)
    gap = Inches(0.25)
    margin = Inches(0.6)
    total_w = int(SLIDE_WIDTH) - int(margin) * 2 - int(gap) * 2
    block_w = Emu(total_w // 3)

    blocks = [
        ("01",
         "Phone & WhatsApp chaos",
         "Orders take 12+ minutes per call. Errors caught after delivery, "
         "blame impossible to assign."),
        ("02",
         "Spreadsheet sprawl",
         "Pricing tiers and promotions fall behind. Margin leaks silently, "
         "rebuilt from memory each quarter."),
        ("03",
         "No paper trail",
         "When something goes wrong, nobody knows who changed what. "
         "Compliance and trust take the hit."),
    ]

    for i, (num, title, body) in enumerate(blocks):
        left = margin + i * (block_w + gap)

        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, block_top, block_w, block_h)
        set_solid_fill(bg, SOFT_BG)
        set_line(bg, HAIRLINE, 0.75)
        bg.adjustments[0] = 0.06

        tb, tf = add_textbox(slide,
                             left + Inches(0.3), block_top + Inches(0.25),
                             block_w - Inches(0.6), block_h - Inches(0.5))
        add_paragraph(tf, num, size_pt=14, bold=True, color=ACCENT,
                      font=FONT_DISPLAY, space_after=8)
        add_paragraph(tf, title, size_pt=18, bold=True, color=INK,
                      font=FONT_DISPLAY, space_after=8)
        add_paragraph(tf, body, size_pt=12, color=MUTED, space_after=0)

    # Closing line
    tb, tf = add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5))
    add_paragraph(tf,
                  "Existing systems treat each channel separately. "
                  "We unify them.",
                  size_pt=14, bold=True, color=INK, space_after=0)

    add_footer(slide, "Nex Order   •   The problem   •   Page 2 of 5")

    add_speaker_notes(slide,
        "Use a call-and-response framing here. Ask the audience: 'Whichever "
        "of these three is biggest for you today, raise a hand.' All three "
        "have the same root cause — fragmented intake. Phone orders, "
        "spreadsheet pricing, and missing audit trails are symptoms of the "
        "same architectural gap. Nex Order is the answer: one entry point, "
        "one rules engine, one log. Bridge to Slide 3: 'Let me show you "
        "how it actually flows.'")


# ---------------------------------------------------------------------------
# Slide 3 — Process Flow (the diagram)
# ---------------------------------------------------------------------------

def build_slide_3(prs: Presentation) -> None:
    slide = add_blank_slide(prs)

    # Eyebrow + headline
    add_accent_bar(slide, left=Inches(0.6), top=Inches(0.6))
    tb, tf = add_textbox(slide, Inches(0.6), Inches(0.78), Inches(6), Inches(0.35))
    add_paragraph(tf, "PROCESS FLOW", size_pt=11, bold=True, color=ACCENT,
                  font=FONT_DISPLAY)

    tb, tf = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.6))
    add_paragraph(tf, "How an order flows through Nex Order.",
                  size_pt=26, bold=True, color=INK, font=FONT_DISPLAY)

    # Layout constants for the three columns
    col_top = Inches(2.1)
    col_h = Inches(4.4)
    col_w_intake = Inches(3.4)
    col_w_pipe = Inches(3.6)
    col_w_fulfill = Inches(3.4)
    gap = Inches(0.4)
    left_intake = Inches(0.6)
    left_pipe = left_intake + col_w_intake + gap
    left_fulfill = left_pipe + col_w_pipe + gap

    # Column headers
    def col_header(left, top, width, label):
        tb, tf = add_textbox(slide, left, top, width, Inches(0.3))
        add_paragraph(tf, label, size_pt=10, bold=True, color=ACCENT,
                      font=FONT_DISPLAY, align=PP_ALIGN.CENTER)

    header_top = col_top - Inches(0.4)
    col_header(left_intake, header_top, col_w_intake, "INTAKE CHANNELS")
    col_header(left_pipe, header_top, col_w_pipe, "UNIFIED PIPELINE")
    col_header(left_fulfill, header_top, col_w_fulfill, "FULFILLMENT")

    # ------ Column A — intake (4 stacked boxes)
    intake_items = [
        ("Customer self-serve", "Shop view on phone or web"),
        ("Office sales rep", "Telesales — call-reference verification"),
        ("Field sales rep", "On-site — customer signature"),
        ("Roadmap: Email / WhatsApp / EDI", "External channels, same gate"),
    ]
    box_h = Emu((int(col_h) - int(Inches(0.3)) * 3) // 4)
    for i, (title, sub) in enumerate(intake_items):
        top = col_top + Emu(i * (int(box_h) + int(Inches(0.1))))
        is_future = i == 3
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left_intake, top, col_w_intake, box_h)
        set_solid_fill(box, WHITE if not is_future else SOFT_BG)
        set_line(box, ACCENT if not is_future else MUTED, 1.25 if not is_future else 0.75)
        box.adjustments[0] = 0.10
        if is_future:
            box.line.dash_style = 7  # dashed

        tb, tf = add_textbox(slide,
                             left_intake + Inches(0.2), top + Inches(0.12),
                             col_w_intake - Inches(0.4), box_h - Inches(0.24))
        add_paragraph(tf, title, size_pt=12, bold=True,
                      color=INK if not is_future else MUTED,
                      font=FONT_DISPLAY, space_after=2)
        add_paragraph(tf, sub, size_pt=9, color=MUTED, space_after=0)

    # ------ Column B — unified pipeline (single tall box)
    pipe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left_pipe, col_top, col_w_pipe, col_h)
    set_solid_fill(pipe, INK)
    set_line(pipe, INK, 1.0)
    pipe.adjustments[0] = 0.06

    tb, tf = add_textbox(slide,
                         left_pipe + Inches(0.25), col_top + Inches(0.25),
                         col_w_pipe - Inches(0.5), col_h - Inches(0.5))
    add_paragraph(tf, "place-order", size_pt=11, bold=True, color=ACCENT,
                  font=FONT_DISPLAY, space_after=4)
    add_paragraph(tf, "Edge Function — single server-side gate",
                  size_pt=12, bold=True, color=WHITE, font=FONT_DISPLAY,
                  space_after=14)

    pipeline_steps = [
        "Auth & role check",
        "Tier pricing resolution",
        "Promotion engine (BOGO, bundles, storewide)",
        "Stock & credit limit checks",
        "Atomic write — orders + items + invoice",
        "Realtime broadcast to all roles",
        "Email confirmation (Resend)",
        "Audit event written",
    ]
    for step in pipeline_steps:
        add_paragraph(tf, "•  " + step, size_pt=11, color=WHITE, space_after=4)

    # ------ Column C — fulfillment (4 stacked boxes)
    fulfill_items = [
        ("Order status pipeline",
         "processing → confirmed → packed → shipped → delivered"),
        ("Invoice lifecycle",
         "pending → issued → paid (auto-email on issue)"),
        ("Stock decrement",
         "Inventory drops; low-stock alert if below threshold"),
        ("Audit log entry",
         "Actor, before/after, optional reason"),
    ]
    for i, (title, sub) in enumerate(fulfill_items):
        top = col_top + Emu(i * (int(box_h) + int(Inches(0.1))))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left_fulfill, top, col_w_fulfill, box_h)
        set_solid_fill(box, WHITE)
        set_line(box, ACCENT, 1.25)
        box.adjustments[0] = 0.10

        tb, tf = add_textbox(slide,
                             left_fulfill + Inches(0.2), top + Inches(0.12),
                             col_w_fulfill - Inches(0.4), box_h - Inches(0.24))
        add_paragraph(tf, title, size_pt=12, bold=True, color=INK,
                      font=FONT_DISPLAY, space_after=2)
        add_paragraph(tf, sub, size_pt=9, color=MUTED, space_after=0)

    # Connector arrows between columns (decorative)
    def connector_arrow(x_start, x_end, y):
        line = slide.shapes.add_connector(1, x_start, y, x_end, y)
        line.line.color.rgb = ACCENT
        line.line.width = Pt(1.5)

    arrow_y = col_top + Emu(int(col_h) // 2)
    connector_arrow(left_intake + col_w_intake, left_pipe, arrow_y)
    connector_arrow(left_pipe + col_w_pipe, left_fulfill, arrow_y)

    # Footer band line
    tb, tf = add_textbox(slide, Inches(0.6), Inches(6.65), Inches(12), Inches(0.4))
    add_paragraph(tf,
                  "One server-side gate. Same rules for every channel. "
                  "Every action logged.",
                  size_pt=12, bold=True, color=INK, align=PP_ALIGN.CENTER)

    add_footer(slide, "Nex Order   •   Process flow   •   Page 3 of 5")

    add_speaker_notes(slide,
        "Walk this slide left to right. Spend 30 seconds per intake channel: "
        "self-serve is the customer's phone or laptop; telesales is the "
        "office rep keying it in with a call-reference number; field reps "
        "capture a customer signature on a tablet; the dashed box is the "
        "roadmap — email-in, WhatsApp, EDI — all designed to enter the same "
        "gate. The middle column is the architectural moat. Every intake "
        "channel hits the same Edge Function. That means pricing, promos, "
        "stock checks, and audit are guaranteed identical. The right column "
        "is what happens next — a deterministic pipeline that ops teams "
        "can rely on.")


# ---------------------------------------------------------------------------
# Slide 4 — Capability grid
# ---------------------------------------------------------------------------

def build_slide_4(prs: Presentation) -> None:
    slide = add_blank_slide(prs)

    add_accent_bar(slide, left=Inches(0.6), top=Inches(0.6))
    tb, tf = add_textbox(slide, Inches(0.6), Inches(0.78), Inches(6), Inches(0.35))
    add_paragraph(tf, "WHAT'S IN THE PLATFORM", size_pt=11, bold=True,
                  color=ACCENT, font=FONT_DISPLAY)

    tb, tf = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.6))
    add_paragraph(tf, "Built for the realities of B2B selling.",
                  size_pt=26, bold=True, color=INK, font=FONT_DISPLAY)

    # 2 rows × 3 columns capability grid
    capabilities = [
        ("Pantry intelligence",
         "Per-customer reorder list, frequency-weighted suggestions, "
         "automatic out-of-stock substitutes."),
        ("Tier pricing & promotions",
         "Gold/Silver/Bronze tiers. BOGO, bundle, and storewide promos "
         "resolved server-side at checkout."),
        ("Routes & scheduled visits",
         "Field reps follow a planned route with stops, change requests, "
         "and on-site signature capture."),
        ("Realtime cross-role visibility",
         "Postgres LISTEN/NOTIFY + TanStack Query invalidation. "
         "Every role sees new orders without refresh."),
        ("Role-based access + audit log",
         "Five roles, RLS-locked tables. Every privileged write recorded "
         "with before/after diff."),
        ("Stock, invoicing & accounting",
         "Live inventory, aging report, credit-limit checks, "
         "low-stock alerts."),
    ]

    grid_top = Inches(2.2)
    grid_h = Inches(4.3)
    margin = Inches(0.6)
    cols = 3
    rows = 2
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)

    total_w = int(SLIDE_WIDTH) - int(margin) * 2 - int(gap_x) * (cols - 1)
    total_h = int(grid_h) - int(gap_y) * (rows - 1)
    cell_w = Emu(total_w // cols)
    cell_h = Emu(total_h // rows)

    for idx, (title, body) in enumerate(capabilities):
        r, c = divmod(idx, cols)
        left = margin + Emu(c * (int(cell_w) + int(gap_x)))
        top = grid_top + Emu(r * (int(cell_h) + int(gap_y)))

        cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left, top, cell_w, cell_h)
        set_solid_fill(cell, WHITE)
        set_line(cell, HAIRLINE, 0.75)
        cell.adjustments[0] = 0.06

        # accent corner bar
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        left + Inches(0.25), top + Inches(0.25),
                                        Inches(0.4), Inches(0.06))
        set_solid_fill(accent, ACCENT)
        set_line(accent, None)

        tb, tf = add_textbox(slide,
                             left + Inches(0.25), top + Inches(0.45),
                             cell_w - Inches(0.5), cell_h - Inches(0.6))
        add_paragraph(tf, title, size_pt=15, bold=True, color=INK,
                      font=FONT_DISPLAY, space_after=6)
        add_paragraph(tf, body, size_pt=11, color=MUTED, space_after=0)

    # Tech footer line
    tb, tf = add_textbox(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4))
    add_paragraph(tf,
                  "13 Edge Functions   •   5 roles   •   Postgres + RLS   •   "
                  "React 19 / TypeScript / Tailwind   •   Vercel + Supabase",
                  size_pt=10, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, "Nex Order   •   Capabilities   •   Page 4 of 5")

    add_speaker_notes(slide,
        "Don't read the grid. Pick three to spotlight — the ones that match "
        "the prospect's biggest pain. For most distributors, lead with "
        "Pantry intelligence (it removes 80% of phone time), then Tier "
        "pricing & promotions (it stops margin leaks), then Routes (it "
        "turns field reps into a controlled channel). Mention Realtime, "
        "Audit, and Stock as 'and you also get…'. Bridge to Slide 5: "
        "'Let me actually show you each of the three I just highlighted.'")


# ---------------------------------------------------------------------------
# Slide 5 — Live demo + next steps
# ---------------------------------------------------------------------------

def build_slide_5(prs: Presentation) -> None:
    slide = add_blank_slide(prs)

    add_accent_bar(slide, left=Inches(0.6), top=Inches(0.6))
    tb, tf = add_textbox(slide, Inches(0.6), Inches(0.78), Inches(6), Inches(0.35))
    add_paragraph(tf, "LIVE DEMO", size_pt=11, bold=True, color=ACCENT,
                  font=FONT_DISPLAY)

    tb, tf = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.0))
    add_paragraph(tf, "Let's see it work.",
                  size_pt=40, bold=True, color=INK, font=FONT_DISPLAY)

    # Left column — what you'll see (3 numbered acts)
    left_col_left = Inches(0.6)
    left_col_top = Inches(2.6)
    left_col_w = Inches(7.5)

    acts = [
        ("01", "A restaurant places an order from their phone",
         "Pantry, promos, tier pricing — no phone call."),
        ("02", "A field rep tops it up during a planned visit",
         "Signature on glass, route updates in realtime."),
        ("03", "The back office sees both orders the moment they land",
         "Confirms, packs, and audits — no refresh required."),
    ]
    row_h = Inches(1.25)
    for i, (num, title, body) in enumerate(acts):
        top = left_col_top + Emu(i * int(row_h))
        # number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        left_col_left, top + Inches(0.1),
                                        Inches(0.7), Inches(0.7))
        set_solid_fill(circle, ACCENT)
        set_line(circle, None)
        ctb, ctf = add_textbox(slide, left_col_left, top + Inches(0.18),
                               Inches(0.7), Inches(0.55))
        add_paragraph(ctf, num, size_pt=16, bold=True, color=WHITE,
                      font=FONT_DISPLAY, align=PP_ALIGN.CENTER)

        tb, tf = add_textbox(slide,
                             left_col_left + Inches(1.0), top + Inches(0.05),
                             left_col_w - Inches(1.0), row_h - Inches(0.1))
        add_paragraph(tf, title, size_pt=15, bold=True, color=INK,
                      font=FONT_DISPLAY, space_after=2)
        add_paragraph(tf, body, size_pt=11, color=MUTED, space_after=0)

    # Right column — Next steps card
    right_left = Inches(8.6)
    right_top = Inches(2.6)
    right_w = Inches(4.1)
    right_h = Inches(3.75)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  right_left, right_top, right_w, right_h)
    set_solid_fill(card, INK)
    set_line(card, INK, 1.0)
    card.adjustments[0] = 0.06

    tb, tf = add_textbox(slide,
                         right_left + Inches(0.3), right_top + Inches(0.3),
                         right_w - Inches(0.6), right_h - Inches(0.6))
    add_paragraph(tf, "NEXT STEPS", size_pt=11, bold=True, color=ACCENT,
                  font=FONT_DISPLAY, space_after=10)
    add_paragraph(tf, "1.  Sandbox account",
                  size_pt=14, bold=True, color=WHITE, space_after=2)
    add_paragraph(tf, "Your team logs in within 24 hours.",
                  size_pt=10, color=HAIRLINE, space_after=12)
    add_paragraph(tf, "2.  Scoping call",
                  size_pt=14, bold=True, color=WHITE, space_after=2)
    add_paragraph(tf, "Map your channels, tiers, and promos.",
                  size_pt=10, color=HAIRLINE, space_after=12)
    add_paragraph(tf, "3.  Rollout plan",
                  size_pt=14, bold=True, color=WHITE, space_after=2)
    add_paragraph(tf, "Pilot with one branch, then scale.",
                  size_pt=10, color=HAIRLINE, space_after=0)

    # Bottom contact strip
    tb, tf = add_textbox(slide, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.4))
    add_paragraph(tf,
                  "nexorder.vercel.app   •   contact: hello@nexorder.app",
                  size_pt=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_footer(slide, "Nex Order   •   Live demo   •   Page 5 of 5")

    add_speaker_notes(slide,
        "Brief slide — under 30 seconds. Set expectations: three acts, one "
        "continuous story. Maya the customer, Charlie the field rep, Alice "
        "the admin. The URL is on screen if anyone wants to follow along on "
        "their laptop. After the demo, return here for next steps and Q&A.")


# ---------------------------------------------------------------------------

def main() -> None:
    prs = make_presentation()
    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH} ({OUTPUT_PATH.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
